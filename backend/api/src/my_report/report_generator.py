from __future__ import annotations

import os
from pathlib import Path
from datetime import datetime
from io import BytesIO
from typing import Iterable, Optional, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.shapes.shapetree import SlideShapes

from .models import TikTokStats, TikTokVideo


# ===============================
# Design System（色・フォント・余白）
# ===============================
class DesignSystem:
    """デザイン定数とカラーパレット"""
    
    # カラーパレット
    BRAND = RGBColor(116, 226, 207)      # #74E2CF - ブランドカラー
    BRAND_DARK = RGBColor(67, 185, 165)  # #43B9A5 - アクセント/強調
    BRAND_LIGHT = RGBColor(169, 240, 228) # #A9F0E4 - 面/帯の淡色
    TEXT = RGBColor(31, 41, 55)          # #1F2937 - 本文
    MUTE = RGBColor(107, 114, 128)       # #6B7280 - 補足/ラベル
    BG = RGBColor(255, 255, 255)         # #FFFFFF - 背景
    PANEL = RGBColor(245, 247, 251)      # #F5F7FB - パネル背景
    BORDER = RGBColor(229, 231, 235)     # #E5E7EB - ボーダー
    ACCENT_ORANGE = RGBColor(255, 140, 0) # #FF8C00 - オレンジ（アイコン用）
    
    # フォント
    FONT_FAMILY = "Meiryo"
    FONT_SIZE_TITLE = Pt(32)
    FONT_SIZE_H1 = Pt(26)
    FONT_SIZE_H2 = Pt(18)
    FONT_SIZE_BODY = Pt(12)
    FONT_SIZE_SMALL = Pt(11)
    
    # レイアウト
    PAGE_WIDTH = Inches(13.33)  # 16:9 WIDE
    PAGE_HEIGHT = Inches(7.5)
    MARGIN = Inches(0.6)
    RADIUS = Inches(0.15)
    GAP = Inches(0.25)


DS = DesignSystem  # 短縮形


# ===============================
# ユーティリティ関数
# ===============================
def _configure_text(paragraph, text: str, *, size: Pt, color: RGBColor, bold: bool = False, align=PP_ALIGN.LEFT):
    """テキストスタイルを統一設定"""
    paragraph.text = text
    paragraph.font.size = size
    paragraph.font.bold = bold
    paragraph.font.name = DS.FONT_FAMILY
    paragraph.font.color.rgb = color
    paragraph.alignment = align


def _add_rect(slide, *, x: float, y: float, w: float, h: float, 
              fill_color: RGBColor, border_color: Optional[RGBColor] = None,
              border_width: Pt = Pt(1), radius: bool = True):
    """角丸矩形を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.width = Pt(0)
    
    shape.shadow.inherit = False
    return shape


def _add_text_box(slide, *, x: float, y: float, w: float, h: float, 
                  text: str, size: Pt, color: RGBColor, bold: bool = False, 
                  align=PP_ALIGN.LEFT):
    """テキストボックスを追加"""
    textbox = slide.shapes.add_textbox(x, y, w, h)
    tf = textbox.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    _configure_text(paragraph, text, size=size, color=color, bold=bold, align=align)
    return textbox


def _format_number(value: Optional[int]) -> str:
    """数値をカンマ区切りにフォーマット"""
    if value is None:
        return "-"
    return f"{value:,}"


# ===============================
# スライドテンプレート
# ===============================
def _create_cover_slide(
    presentation: Presentation,
    *,
    account_name: str,
    period_label: str,
    start_date: datetime,
    end_date: datetime,
    generated_at: datetime,
) -> None:
    """カバースライド"""
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    
    width = presentation.slide_width
    height = presentation.slide_height
    
    # 背景
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DS.BG
    
    # === ヘッダー画像（斜めグラデーション）を配置 ===
    header_height = height * 0.66  # 2/3の高さ
    
    # 画像パスを取得（このファイルと同じディレクトリのimagesフォルダ内）
    current_dir = Path(__file__).parent
    header_image_path = current_dir / "images" / "teal_header_block_MED.png"
    
    if header_image_path.exists():
        # 画像を配置（幅いっぱい、上から2/3の高さ）
        slide.shapes.add_picture(
            str(header_image_path),
            0, 0,  # 左上から
            width, header_height  # 幅いっぱい、2/3の高さ
        )
    else:
        # 画像がない場合はフォールバック（単色の帯）
        _add_rect(slide, x=0, y=0, w=width, h=header_height,
                  fill_color=DS.BRAND, radius=False)
    
    # === タイトルエリア（画像の上に配置） ===
    title_y = header_height * 0.35  # 帯の中央よりやや上
    
    # 英語タイトル
    _add_text_box(slide, 
                  x=DS.MARGIN, y=title_y, 
                  w=width - DS.MARGIN * 2, h=Inches(0.4),
                  text="TikTok Monthly Report", 
                  size=Pt(20), color=DS.BG)
    
    # 日本語タイトル（大きく）
    _add_text_box(slide, 
                  x=DS.MARGIN, y=title_y + Inches(0.45), 
                  w=width - DS.MARGIN * 2, h=Inches(0.8),
                  text="TikTok月次レポート", 
                  size=Pt(42), color=DS.BG, bold=True)
    
    # === 下部の白いカード（画像の下） ===
    card_top = header_height + Inches(0.15)
    card_height = height - card_top - DS.MARGIN
    
    _add_rect(slide, 
              x=DS.MARGIN, y=card_top,
              w=width - DS.MARGIN * 2, h=card_height,
              fill_color=DS.BG, border_color=DS.BORDER)
    
    # アカウント情報
    info_x = DS.MARGIN + Inches(0.4)
    info_y = card_top + Inches(0.3)
    
    display_account = account_name if account_name.startswith('@') else f"@{account_name}"
    _add_text_box(slide, x=info_x, y=info_y, 
                  w=width - DS.MARGIN * 2 - Inches(0.8), h=Inches(0.5),
                  text=display_account, size=Pt(28), 
                  color=DS.TEXT, bold=True)
    
    # データ期間と作成日（●をオレンジ、テキストを黒に）
    info_y += Inches(0.6)
    
    # データ期間
    period_box = slide.shapes.add_textbox(info_x, info_y, 
                                          width - DS.MARGIN * 2 - Inches(0.8), 
                                          Inches(0.4))
    period_tf = period_box.text_frame
    period_para = period_tf.paragraphs[0]
    
    # ●（オレンジ）
    period_para.text = "● "
    period_para.font.size = Pt(15)
    period_para.font.name = DS.FONT_FAMILY
    period_para.font.color.rgb = DS.ACCENT_ORANGE
    
    # テキスト部分（黒）
    period_run = period_para.add_run()
    period_run.text = f"データの集計期間  ……………  {start_date:%Y/%m/%d} – {end_date:%Y/%m/%d} ({period_label})"
    period_run.font.size = Pt(15)
    period_run.font.name = DS.FONT_FAMILY
    period_run.font.color.rgb = DS.TEXT
    
    # 作成日
    info_y += Inches(0.35)
    created_box = slide.shapes.add_textbox(info_x, info_y, 
                                           width - DS.MARGIN * 2 - Inches(0.8), 
                                           Inches(0.4))
    created_tf = created_box.text_frame
    created_para = created_tf.paragraphs[0]
    
    # ●（オレンジ）
    created_para.text = "● "
    created_para.font.size = Pt(15)
    created_para.font.name = DS.FONT_FAMILY
    created_para.font.color.rgb = DS.ACCENT_ORANGE
    
    # テキスト部分（黒）
    created_run = created_para.add_run()
    created_run.text = f"作成日  …………………………………  {generated_at:%Y/%m/%d}"
    created_run.font.size = Pt(15)
    created_run.font.name = DS.FONT_FAMILY
    created_run.font.color.rgb = DS.TEXT


def _create_section_slide(presentation: Presentation, *, title: str, notes: str) -> None:
    """セクション扉スライド"""
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    
    width = presentation.slide_width
    
    # 背景
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DS.BG
    
    # サブタイトル
    _add_text_box(slide, x=DS.MARGIN, y=Inches(0.42), 
                  w=width - DS.MARGIN * 2, h=Inches(0.3),
                  text="TikTok Monthly Report", size=DS.FONT_SIZE_BODY, 
                  color=DS.MUTE)
    
    # メインタイトル
    _add_text_box(slide, x=DS.MARGIN, y=Inches(1.10), 
                  w=width - DS.MARGIN * 2, h=Inches(1.0),
                  text=title, size=Pt(36), color=DS.BRAND, bold=True)
    
    # 下部の淡い帯
    _add_rect(slide, x=0, y=Inches(4.60), w=width, h=Inches(1.60),
              fill_color=DS.BRAND_LIGHT, radius=False)
    
    # 注釈カード
    _add_rect(slide, 
              x=DS.MARGIN, y=Inches(4.65),
              w=width - DS.MARGIN * 2, h=Inches(1.50),
              fill_color=DS.BG, border_color=DS.BORDER)
    
    _add_text_box(slide, 
                  x=DS.MARGIN + Inches(0.3), y=Inches(4.85),
                  w=width - DS.MARGIN * 2 - Inches(0.6), h=Inches(1.10),
                  text=notes, size=DS.FONT_SIZE_BODY, color=DS.TEXT)


def _create_summary_slide(presentation: Presentation, stats: TikTokStats) -> None:
    """サマリースライド（KPIカード）"""
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    
    width = presentation.slide_width
    
    # 背景
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DS.PANEL
    
    # タイトル
    _add_text_box(slide, x=DS.MARGIN, y=Inches(0.40), 
                  w=width - DS.MARGIN * 2, h=Inches(0.5),
                  text="総括サマリー", size=DS.FONT_SIZE_H1, 
                  color=DS.BRAND, bold=True)
    
    # KPIカード
    cards_top = Inches(1.2)
    card_width = Inches(3.0)
    card_height = Inches(1.8)
    
    kpis = [
        ("視聴増分", _format_number(stats.viewGrowth), "Views"),
        ("平均視聴/動画", _format_number(stats.avgViewCount), "Avg Views"),
        ("フォロワー総数", _format_number(stats.followerCount), "Followers"),
        ("フォロワー増加", _format_number(stats.followerGrowth), "Net Add"),
    ]
    
    for idx, (title, value, label) in enumerate(kpis):
        left = DS.MARGIN + (card_width + DS.GAP) * idx
        
        # カード背景
        _add_rect(slide, x=left, y=cards_top, w=card_width, h=card_height,
                  fill_color=DS.BG, border_color=DS.BORDER)
        
        # タイトル
        _add_text_box(slide, x=left + Inches(0.2), y=cards_top + Inches(0.2),
                      w=card_width - Inches(0.4), h=Inches(0.3),
                      text=title, size=Pt(13), color=DS.TEXT)
        
        # 値
        _add_text_box(slide, x=left + Inches(0.2), y=cards_top + Inches(0.6),
                      w=card_width - Inches(0.4), h=Inches(0.8),
                      text=value, size=Pt(28), color=DS.BRAND, bold=True)
        
        # ラベル
        _add_text_box(slide, x=left + Inches(0.2), y=cards_top + Inches(1.45),
                      w=card_width - Inches(0.4), h=Inches(0.3),
                      text=label, size=DS.FONT_SIZE_SMALL, color=DS.MUTE)


def _create_top_videos_slide(presentation: Presentation, videos: Sequence[TikTokVideo]) -> None:
    """トップ動画スライド"""
    top_videos = sorted(videos, key=lambda v: v.viewCount, reverse=True)[:6]
    if not top_videos:
        return
    
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    width = presentation.slide_width
    
    # 背景
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DS.PANEL
    
    # タイトル
    _add_text_box(slide, x=DS.MARGIN, y=Inches(0.40), 
                  w=width - DS.MARGIN * 2, h=Inches(0.5),
                  text="視聴トップ動画", size=DS.FONT_SIZE_H1, 
                  color=DS.BRAND, bold=True)
    
    # グリッド配置
    card_width = Inches(4.0)
    card_height = Inches(1.6)
    start_top = Inches(1.2)
    
    for idx, video in enumerate(top_videos):
        row = idx // 2
        col = idx % 2
        left = DS.MARGIN + col * (card_width + DS.GAP)
        top = start_top + row * (card_height + DS.GAP)
        
        # カード背景
        _add_rect(slide, x=left, y=top, w=card_width, h=card_height,
                  fill_color=DS.BG, border_color=DS.BORDER)
        
        # タイトル
        title_text = f"#{idx + 1} {video.title or '(タイトルなし)'}"
        _add_text_box(slide, x=left + Inches(0.2), y=top + Inches(0.15),
                      w=card_width - Inches(0.4), h=Inches(0.35),
                      text=title_text, size=Pt(13), color=DS.TEXT, bold=True)
        
        # メタ情報
        meta_text = f"投稿日: {video.createTime[:10]} / 視聴数: {_format_number(video.viewCount)}"
        _add_text_box(slide, x=left + Inches(0.2), y=top + Inches(0.55),
                      w=card_width - Inches(0.4), h=Inches(0.3),
                      text=meta_text, size=DS.FONT_SIZE_SMALL, color=DS.TEXT)
        
        # エンゲージメント
        eng_text = f"いいね: {_format_number(video.likeCount)}  コメント: {_format_number(video.commentCount)}  シェア: {_format_number(video.shareCount)}"
        _add_text_box(slide, x=left + Inches(0.2), y=top + Inches(0.90),
                      w=card_width - Inches(0.4), h=Inches(0.5),
                      text=eng_text, size=Pt(10), color=DS.MUTE)


# ===============================
# メイン関数
# ===============================
def build_tiktok_report_presentation(
    *,
    stats: TikTokStats,
    videos: Sequence[TikTokVideo],
    account_name: str,
    period_label: str,
    start_date: datetime,
    end_date: datetime,
    generated_at: Optional[datetime] = None,
) -> BytesIO:
    """TikTokレポートのPowerPointを生成"""
    generated_at = generated_at or datetime.now()
    
    presentation = Presentation()
    presentation.slide_width = DS.PAGE_WIDTH
    presentation.slide_height = DS.PAGE_HEIGHT
    
    # カバー
    _create_cover_slide(
        presentation,
        account_name=account_name,
        period_label=period_label,
        start_date=start_date,
        end_date=end_date,
        generated_at=generated_at,
    )
    
    # サマリー
    _create_summary_slide(presentation, stats)
    
    # セクション扉
    _create_section_slide(
        presentation, 
        title="コンテンツ分析", 
        notes="POINT\n- 視聴トップ動画を確認\n- エンゲージメントの高いコンテンツを把握"
    )
    
    # トップ動画
    _create_top_videos_slide(presentation, videos)
    
    stream = BytesIO()
    presentation.save(stream)
    stream.seek(0)
    return stream
